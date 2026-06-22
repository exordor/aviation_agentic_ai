"""Build the ATCSCC master-project defense slide deck (reproducible, tracked).

Generates reports/final/atcscc_defense_deck.pptx from the current ATCSCC thesis
evidence (report + RQ3 head-to-head). Run:

    uv run python scripts/build_atcscc_defense_deck.py

This replaces the PHAK-era defense deck with current ATCSCC content.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "reports" / "final" / "atcscc_defense_deck.pptx"
ARCH = ROOT / "reports" / "final" / "assets" / "atcscc_architecture.svg"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = "0F172A"
MUTED = "475569"
ACCENT = "B91C1C"
LIGHT_BG = "F8FAFC"


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(INK)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(MUTED)


def _add_content_slide(prs: Presentation, title: str, bullets: list[tuple[str, int]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(INK)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True
    for text, level in bullets:
        para = tf.paragraphs[0] if tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs else tf.add_paragraph()
        para.text = text
        para.level = level
        para.font.size = Pt(20 if level == 0 else 16)
        para.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(
            INK if level == 0 else MUTED
        )
        para.space_after = Pt(6)


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(INK)

    cols = len(headers)
    rcount = len(rows) + 1
    tbl_shape = slide.shapes.add_table(rcount, cols, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5 * rcount))
    table = tbl_shape.table
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htext
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(15)
                para.font.bold = (j == 0)


def _add_image_slide(prs: Presentation, title: str, note: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    nb = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.4))
    ntf = nb.text_frame
    ntf.word_wrap = True
    p = ntf.paragraphs[0]
    p.text = note
    p.font.size = Pt(16)
    p.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(MUTED)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(
        prs,
        "Schema-Constrained Agentic KG-RAG",
        "Evidence-Grounded QA over FAA ATCSCC Advisories  ·  Master Project Defense",
    )

    _add_content_slide(
        prs,
        "Motivation & Problem",
        [
            ("ATCSCC advisories = public, semi-structured NAS traffic-management notices.", 0),
            ("Short texts with identifiers, affected NAS elements, time windows, causes.", 1),
            ("Facts are source-visible and checkable against evidence spans.", 1),
            ("Problem: can a lightweight schema constrain LLM extraction, support agentic", 0),
            ("validation/refinement, and yield an event graph that improves grounded QA?", 1),
            ("Contribution = method integration under ONE bounded source family.", 0),
            ("NOT a complete aviation ontology; NOT operational ATC decision support.", 1),
        ],
    )

    _add_content_slide(
        prs,
        "Research Questions",
        [
            ("RQ1 — Schema-constrained extraction: valid, evidence-linked event records?", 0),
            ("RQ2 — Agentic validation-refinement: reduce schema violations & unsupported relations?", 0),
            ("RQ3 — KG-RAG grounding: improve evidence grounding & citations vs vector-only?", 0),
            ("RQ4 — Failure boundary: what failure types remain; where is human review needed?", 0),
            ("Evaluation philosophy: layered metrics, NO mixed overall score.", 0),
        ],
    )

    _add_image_slide(
        prs,
        "Method: Five-Block Architecture",
        "See reports/final/assets/atcscc_architecture.svg (PNG render recommended for the deck).\n"
        "advisory → schema → agentic extraction (S0–S4 + validator/refiner/critic) → "
        "evidence-linked event graph → vector/graph/hybrid/routed KG-RAG → answers + citations.\n\n"
        "Reproduce end-to-end:  uv run aviation-ai demo   (offline, no API key)",
    )

    _add_content_slide(
        prs,
        "Data & Schema",
        [
            ("Retrospective ATCSCC snapshot (2026-05-14 → 2026-05-20).", 0),
            ("867 advisory pages → 100 reviewed gold records.", 1),
            ("Classes: GroundDelay(16) · GroundStop(21) · ReRoute(23) · TrafficMgmt(40).", 1),
            ("Schema = NASA ATMONTO-derived ATCSCC slice (not the full ontology).", 0),
            ("18 classes · 22 properties · evidence-span required · profile-gap policy.", 1),
            ("S0 owns deterministic fields; S3/S4 may add but not overwrite semantic fields.", 1),
        ],
    )

    _add_table_slide(
        prs,
        "RQ1 — Schema-Constrained Extraction (100 reviewed records, strict semantic F1)",
        ["System", "Precision", "Recall", "F1", "Schema violation", "Structural accept"],
        [
            ["S0 rule-only", "0.816", "0.710", "0.759", "0.078", "0.922"],
            ["S1b LLM canonicalized", "0.497", "0.144", "0.224", "0.584", "0.416"],
            ["S2 schema-slice LLM", "0.206", "0.184", "0.195", "0.175", "0.825"],
            ["S3 validator-repair", "0.242", "0.132", "0.171", "0.104", "0.897"],
            ["S4 hybrid enrichment", "0.686", "0.774", "0.727", "0.000", "1.000"],
        ],
    )

    _add_content_slide(
        prs,
        "RQ1 Reading + RQ2 Agentic Loop",
        [
            ("RQ1: deterministic S0 (0.759) and S4 hybrid (0.727) are strongest.", 0),
            ("Pure-LLM arms (S1b/S2/S3) markedly weaker — negative result, reported directly.", 1),
            ("Schema constraint effective: S1b structural accept 0.42 vs S2/S4 ≥ 0.82.", 1),
            ("Provenance completeness 1.0 · evidence-in-source 1.0 · 448 valid triples.", 1),
            ("RQ2: validator/refiner/critic loop records auditable repair/reject decisions.", 0),
            ("S3 repair success 0.897 · S4 repair success 1.0 · quarantine 0.", 1),
            ("Loop = diagnostic/repair framework; not autonomous ontology construction.", 1),
        ],
    )

    _add_table_slide(
        prs,
        "RQ3 — Head-to-Head: KG-RAG vs Vector-Only (same 30 ATCSCC questions, gpt-5.4-mini)",
        ["Mode", "Questions", "Correctness", "Unsupported claim", "Citation recall"],
        [
            ["KG-RAG (routed, graph+critic)", "30", "0.967", "0.017", "0.608"],
            ["Vector-only (live tfidf, no graph)", "30", "0.500", "0.500", "0.372"],
        ],
    )

    _add_content_slide(
        prs,
        "RQ3 Reading",
        [
            ("KG-RAG ~doubles correctness; ~30x lower unsupported-claim rate.", 0),
            ("Gain concentrated on entity/cause/status templates:", 0),
            ("AFFECTED-NAS-ELEMENTS, CAUSE-CONDITION, STATUS-ACTION.", 1),
            ("Vector-only ties/wins on TIME-WINDOW, ABSTENTION, ROUTE-SEMANTICS.", 1),
            ("Recall@5 is non-discriminating on ATCSCC (0.685 for most modes).", 0),
            ("The graph's value is at the answer-correctness layer, not retrieval recall.", 1),
            ("Forcing graph everywhere hurts abstention (F1 0.52) — routing matters.", 1),
            ("STILL: source-bounded, 30 questions — not a universal GraphRAG claim.", 0),
        ],
    )

    _add_content_slide(
        prs,
        "RQ4 — Failure & Review Boundary",
        [
            ("9 generated-answer cases packaged for review (incl. 3 failures).", 0),
            ("3 failures adjudicated as profile/gold-boundary cases.", 1),
            ("Profile-decision what-if: predicate whitelist would fix those 3,", 1),
            ("but does NOT replace strict metrics or change gold/profile artifacts.", 1),
            ("Failure taxonomy: extraction · retrieval · profile/gold-boundary ·", 0),
            ("answer-overreach · human-review cases.", 1),
            ("Automated diagnostics ≠ human/expert review. Operational use out of scope.", 0),
        ],
    )

    _add_content_slide(
        prs,
        "Claim Safety",
        [
            ("CAN claim: schema constrains extraction; provenance preserved; KG-RAG", 0),
            ("improves some source-bounded grounding diagnostics; failures categorized.", 1),
            ("MUST NOT claim: complete ontology; universal GraphRAG superiority;", 0),
            ("expert certification; operational ATC/flight safety.", 1),
            ("Reviewer-defense + SOTA-goal audits both pass for internal/diagnostic claims.", 0),
        ],
    )

    _add_content_slide(
        prs,
        "Reproducibility & Demo",
        [
            ("CLI end-to-end demo (offline):  uv run aviation-ai demo", 0),
            ("RQ3 vector-only arm:  scripts/build_nasa_atmonto_s7_llm_answer_generation.py", 0),
            ("--modes vector-only --run-llm", 1),
            ("Quality gates:  uv run ruff check .   &&   uv run pytest -q", 0),
            ("Full report:  reports/final/atcscc_thesis_report.md", 0),
            ("Dashboard:  reports/stages/thesis_experiment_dashboard.md", 0),
        ],
    )

    _add_content_slide(
        prs,
        "Conclusion",
        [
            ("Reproducible schema-constrained Agentic KG-RAG over retrospective ATCSCC advisories.", 0),
            ("Schema constrains extraction & supports deterministic validation.", 0),
            ("Agentic loop gives inspectable repair/rejection signals.", 0),
            ("Matched head-to-head: KG-RAG improves correctness & cuts unsupported claims", 0),
            ("on relation-oriented questions vs vector-only baseline.", 1),
            ("Bounded method for evidence-grounded advisory QA —", 0),
            ("not a certified ontology or live operational system.", 1),
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(ROOT)} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
